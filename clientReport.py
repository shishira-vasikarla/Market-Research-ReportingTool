from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE

import streamlit as st
import pandas as pd
import openai
from dotenv import load_dotenv
import os
import uuid
import tempfile

# Load environment variables
load_dotenv()

# Set OpenAI API key
openai.api_key = os.getenv('OPENAI_API_KEY')

st.title('Interactive Q&A with Uploaded Files and Reporting')

# File uploaders
uploaded_company_csv = st.file_uploader("Choose a CSV file containing company information ('company URL' and 'Information about the company.')", type=['csv'], key="company_csv")
uploaded_questions_csv = st.file_uploader("Choose a csv file containing questions ('Questions.')", type=['csv'], key="questions_csv")
uploaded_report_questions_csv = st.file_uploader("Choose a CSV file for report questions ('Report Question number' and 'Report Questions')", type=['csv'], key="report_questions_csv")

# Initialize or retrieve the Q&A history
if 'qna_history' not in st.session_state:
    st.session_state.qna_history = []
    


def download_qna_responses():
    if 'qna_history' in st.session_state and st.session_state.qna_history:
        # Convert Q&A history to DataFrame
        qna_df = pd.DataFrame(st.session_state.qna_history)
        # Create a temporary file to save the CSV
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.csv')
        qna_df.to_csv(temp_file.name, index=False)

        # Open the file and offer it for download
        with open(temp_file.name, "rb") as file:
            st.download_button(label="Download Q&A Responses as CSV", data=file, file_name="QnAResponses.csv", mime="text/csv")

# Button to download the Q&A responses in CSV format
if st.button("Download Q&A Responses"):
    download_qna_responses()

def generate_answer_with_chatgpt(question, company_info):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful market and business analyst."},
                {"role": "user", "content": f"Company Information: {company_info}\nQuestion: {question}"}
            ],
        )
        return response.choices[0].message['content'].strip()
    except Exception as e:
        st.error(f"An error occurred: {str(e)}")
        return None

def process_files():
    if uploaded_company_csv and uploaded_questions_csv:
        company_df = pd.read_csv(uploaded_company_csv)
        questions_df = pd.read_csv(uploaded_questions_csv)
        
        for _, question_row in questions_df.iterrows():
            question = question_row['Questions.']
            for _, company_row in company_df.iterrows():
                company_info = company_row['Information about the company.']
                answer = generate_answer_with_chatgpt(question, company_info)
                if answer:
                    st.session_state.qna_history.append({
                        "ID": str(uuid.uuid4()),
                        "CompanyURL": company_row['company URL'],
                        "Question": question,
                        "Answer": answer
                    })

if st.button("Generate Q&A Responses", key="generate_qna_responses"):
    process_files()


if st.session_state.qna_history:
    for item in st.session_state.qna_history:
        st.write(f"Company URL: {item['CompanyURL']}")
        st.write(f"Question: {item['Question']}")
        st.write(f"Answer: {item['Answer']}")
        st.markdown("---")

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib import colors
import pandas as pd
import markdown2 
import re

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import pandas as pd
import re

from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER

from reportlab.platypus import Paragraph, Spacer

def markdown_to_paragraph(text, styles):
    """
    Process text to identify and format headings, bullet points, and ensure text uniformity.
    Returns a list of Paragraph objects with appropriate formatting.
    """
    # Helper function to create a styled paragraph
    def create_paragraph(content, style):
        return Paragraph(content, style)
    
    # Split the text into lines for processing
    lines = text.split('\n')
    elements = []

    for line in lines:
        if line.startswith('### ') or line.startswith('#### '):
            # Remove markdown heading indicators
            clean_line = line.replace('### ', '').replace('#### ', '')
            # Use the Heading2 style for uniformity, or create a custom style if needed
            heading_style = styles['Heading2']
            heading_style.alignment = TA_LEFT
            elements.append(create_paragraph('<b>{}</b>'.format(clean_line), heading_style))
        elif line.strip():  # Non-empty line
            # Regular text line
            body_text_style = styles['BodyText']
            body_text_style.firstLineIndent = 0  # Adjust indent if needed
            # Remove additional markdown symbols like **, - **, and #
            # Updated to handle double asterisks surrounding words more robustly
            processed_line = line.replace('- **', '').replace('#', '')
            processed_line = ' '.join([word.replace('**', '') if '**' in word else word for word in processed_line.split()])
            elements.append(create_paragraph(processed_line, body_text_style))
        else:  # Empty line, insert a spacer
            elements.append(Spacer(1, 12))

    return elements


def convert_csv_to_pdf(csv_file_path, pdf_file_path):
    styles = getSampleStyleSheet()

    # Check and adjust existing styles or define new if necessary
    if 'Heading3' not in styles:
        styles.add(ParagraphStyle(name='Heading3', parent=styles['Heading1'], fontSize=14, leading=18, spaceAfter=6))
    if 'Heading4' not in styles:
        styles.add(ParagraphStyle(name='Heading4', parent=styles['Heading2'], fontSize=12, leading=16, spaceAfter=6))

    report = SimpleDocTemplate(pdf_file_path, pagesize=letter)
    elements = []

    # Title for the PDF report
    elements.append(Paragraph('Report Analysis', styles['Title']))
    elements.append(Spacer(1, 12))

    # Read CSV and generate report content
    data_frame = pd.read_csv(csv_file_path)
    for index, row in data_frame.iterrows():
        answer = row['Report Answer']  # Column name in your CSV
        for paragraph in markdown_to_paragraph(answer, styles):
            elements.append(paragraph)
            elements.append(Spacer(1, 12))

    # Build the PDF with the collected elements
    report.build(elements)



# Function to generate CSV and PDF report, then return their file paths
def generate_report_and_files():
    if uploaded_report_questions_csv is not None and 'qna_history' in st.session_state:
        report_questions_df = pd.read_csv(uploaded_report_questions_csv)
        report_responses = []

        for _, rq_row in report_questions_df.iterrows():
            rq_question = rq_row['Report Questions']
            prompt = f"Analyze the following information and answer the question: {rq_question}\n\n"
            
            for item in st.session_state.qna_history:
                prompt += f"Company URL: {item['CompanyURL']}, Question: {item['Question']}, Answer: {item['Answer']}\n"
            
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are a knowledgeable analyst providing insights based on the given data."},
                        {"role": "user", "content": prompt}
                    ],
                )
                report_response = response.choices[0].message['content'].strip()
                report_responses.append({"Report Question": rq_question, "Report Answer": report_response})
            except Exception as e:
                st.error(f"An error occurred while generating the report: {str(e)}")
                return None, None

        # Create DataFrame from responses and save as CSV
        report_df = pd.DataFrame(report_responses)
        temp_csv = tempfile.NamedTemporaryFile(delete=False, suffix='.csv')
        csv_file_path = temp_csv.name
        report_df.to_csv(csv_file_path, index=False)

        # Convert CSV to PDF
        pdf_file_path = csv_file_path.replace('.csv', '.pdf')
        convert_csv_to_pdf(csv_file_path, pdf_file_path)

        return csv_file_path, pdf_file_path

    return None, None

# Button to generate and download report in CSV format
if st.button("Generate and Download Report in CSV", key="generate_download_report_csv"):
    csv_file_path, _ = generate_report_and_files()
    if csv_file_path:
        with open(csv_file_path, "rb") as file:
            st.download_button(label="Download Report as CSV", data=file, file_name="ReportGenerated.csv", mime="text/csv")

# Button to generate and download report in PDF format
if st.button("Generate and Download Report in PDF", key="generate_download_report_pdf"):
    _, pdf_file_path = generate_report_and_files()
    if pdf_file_path:
        with open(pdf_file_path, "rb") as file:
            st.download_button(label="Download Report as PDF", data=file, file_name="ReportGenerated.pdf", mime="application/pdf")



def split_content(content, word_limit=175):
    """
    Splits the content into chunks where each chunk has up to `word_limit` words.
    """
    words = content.split()
    for i in range(0, len(words), word_limit):
        yield ' '.join(words[i:i+word_limit])

def process_content_for_slides(content, word_limit=170):
    """
    Splits content into chunks based on word limit, and identifies headings and bullet points.
    Returns a list of tuples: (is_heading, text) where `is_heading` is a boolean.
    """
    words = content.split()
    chunks = []
    current_chunk = []
    current_word_count = 0

    for word in words:
        if current_word_count + len(word.split()) <= word_limit:
            current_chunk.append(word)
            current_word_count += len(word.split())
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_word_count = len(word.split())

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    processed_chunks = []
    for chunk in chunks:
        lines = chunk.split('\n')
        processed_lines = []
        for line in lines:
            if line.startswith('### '):
                # Treat as a heading
                processed_lines.append((True, line[4:]))
            else:
                # Regular text or bullet point
                processed_lines.append((False, line))
        processed_chunks.append(processed_lines)

    return processed_chunks

def add_content_slide(prs, content):
    """
    Adds slides with content to the presentation. Headings are bolded, and the rest of the text is standard.
    """
    processed_content = process_content_for_slides(content)

    for slide_content in processed_content:
        slide_layout = prs.slide_layouts[5]  # Use a blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Define text box size and position
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True  # Ensure text wraps within the textbox

        for is_heading, text in slide_content:
            p = tf.add_paragraph()
            p.text = text

            # Check if this line is a heading and should be bolded
            if is_heading:
                run = p.add_run()
                run.font.bold = True
                p.font.size = Pt(18)  # Optional: You might want to use a larger font for headings
            else:
                # This ensures non-heading text is not bolded
                p.font.bold = False
                p.font.size = Pt(16)  # Standard text size
                
            p.level = 0  # Adjust indentation level if needed (e.g., for bullet points)

# Continue with the rest of your csv_to_ppt_enhanced function and other logic



def csv_to_ppt_enhanced(csv_file_path, ppt_file_path):
    prs = Presentation()
    df = pd.read_csv(csv_file_path)

    for index, row in df.iterrows():
        # Use only 'Report Answer' as content, exclude 'Report Question'
        content = row['Report Answer'].replace('**', '')  # Convert '**' to newline or handle accordingly

        # Add content to slides, splitting as necessary
        add_content_slide(prs, content)

    prs.save(ppt_file_path)

# Example usage within Streamlit upon button click
if st.button("Generate and Download Enhanced Report in PPT", key="generate_download_enhanced_report_ppt"):
    if 'qna_history' in st.session_state and uploaded_report_questions_csv is not None:
        csv_file_path, _ = generate_report_and_files()
        ppt_file_path = csv_file_path.replace('.csv', '_enhanced.pptx')
        csv_to_ppt_enhanced(csv_file_path, ppt_file_path)
        
        with open(ppt_file_path, "rb") as file:
            st.download_button(label="Download Enhanced Report as PPT", data=file, file_name="ReportGenerated_enhanced.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

